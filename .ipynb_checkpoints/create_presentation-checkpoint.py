from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Cement Demand Forecasting Across Multiple Sites"
subtitle.text = "Machine Learning Solution for Construction Industry\nInventory Optimization & Demand Prediction"

# Slide 2: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Project Overview"
content.text = """• Machine learning solution for forecasting cement demand
• Analyzes historical operational data and weather patterns
• Predicts future cement consumption across multiple sites
• Optimizes inventory management and reduces costs
• Integrates site characteristics and behavioral patterns"""

# Slide 3: Business Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Business Problem"
content.text = """Construction sites face critical challenges:

• Variable demand patterns across different sites
• Weather-dependent construction activities
• Different site behaviors (aggressive vs conservative)
• Storage capacity constraints
• Cost implications of stockouts vs excess inventory

Solution: Predictive analytics for optimal inventory management"""

# Slide 4: Dataset Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Dataset Description"
content.text = """SQLite Database with 3 main tables:

Sites Table:
• Site ID, Region (North/South/East)
• Silo capacity, Ordering behavior

Operations Table:
• Daily transactions, Planned vs actual consumption
• Inventory levels, Deliveries, Weather data

CementTypes Table:
• CEM_I, CEM_II, CEM_III specifications"""

# Slide 5: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Key Features"
content.text = """🏗️ Multi-site Analysis
   Forecasting across different geographic regions

🌦️ Weather Integration
   Incorporates rainfall and temperature data

📊 Behavioral Patterns
   Accounts for site-specific ordering behaviors

📦 Inventory Optimization
   Balances demand forecasting with storage constraints

📈 Performance Tracking
   Model evaluation and KPI monitoring"""

# Slide 6: Strategic Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Strategic Approach"
content.text = """📊 Data-Driven Decision Making:
• Historical pattern analysis for trend identification
• Real-time monitoring and adaptive forecasting

🎯 Multi-Level Forecasting Strategy:
• Site-specific models for local patterns
• Regional aggregation for supply chain optimization
• Cement type differentiation for specialized demand

⚡ Agile Implementation:
• Phased rollout across sites
• Continuous model improvement
• Stakeholder feedback integration"""

# Slide 7: Technical Implementation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Technical Implementation"
content.text = """Data Processing:
• SQLite database integration
• Feature engineering and data cleaning
• Time-based feature creation

Machine Learning Pipeline:
• Decision Trees, Time Series models
• Historical consumption + weather features
• Time-series cross-validation

Key Technologies:
• Python, Pandas, Scikit-learn
• Jupyter Notebook, SQLite"""

# Slide 8: Model Performance
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Model Performance & Metrics"
content.text = """Forecast Accuracy:
• MAPE (Mean Absolute Percentage Error)
• RMSE (Root Mean Square Error)
• MAE (Mean Absolute Error)

Inventory KPIs:
• Stockout frequency reduction
• Excess inventory cost optimization

Operational Metrics:
• Delivery optimization
• Capacity utilization improvement"""

# Slide 9: Implementation Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Implementation Strategy"
content.text = """Phase 1: Foundation (Months 1-2)
• Data infrastructure setup
• Historical data analysis and cleaning
• Baseline model development

Phase 2: Pilot Testing (Months 3-4)
• Deploy models at 2-3 selected sites
• Monitor performance and gather feedback
• Refine algorithms based on real-world results

Phase 3: Full Deployment (Months 5-6)
• Roll out across all sites
• Integration with existing systems
• Staff training and change management"""

# Slide 10: Risk Mitigation Strategies
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Risk Mitigation Strategies"
content.text = """🛡️ Data Quality Assurance:
• Automated data validation checks
• Outlier detection and handling protocols
• Regular data audit procedures

⚠️ Model Reliability:
• Ensemble methods for robust predictions
• Confidence intervals for uncertainty quantification
• Fallback to historical averages when needed

🔄 Operational Continuity:
• Manual override capabilities
• Regular model retraining schedules
• Performance monitoring dashboards"""

# Slide 11: Results & Business Impact
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results & Business Impact"
content.text = """Key Insights:
• Weather and historical consumption are key predictors
• Site behavior patterns significantly impact demand variability
• Regional differences require tailored forecasting approaches

Business Benefits:
✅ Optimized inventory levels across sites
✅ Reduced stockout incidents
✅ Improved delivery scheduling
✅ Enhanced capacity utilization
✅ Cost reduction through better planning"""

# Slide 12: Success Metrics & KPIs
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Success Metrics & KPIs"
content.text = """📈 Forecast Accuracy Targets:
• MAPE < 15% for weekly forecasts
• RMSE reduction of 25% vs baseline
• 90% prediction confidence intervals

💰 Financial Impact Goals:
• 20% reduction in inventory holding costs
• 30% decrease in stockout incidents
• 15% improvement in delivery efficiency

⏱️ Operational Excellence:
• 95% system uptime
• <2 hour model refresh cycles
• 100% site coverage within 6 months"""

# Slide 13: Project Structure
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Project Structure"
content.text = """📁 Main Components:
• Cement_Demand_Forecasting.ipynb - Analysis notebook
• MIG_Cement_Records.db - SQLite database
• cement_forecast_results.parquet - Predictions
• Model_performance_summary.csv - Metrics
• feature_importance_dt.png - Visualizations

📋 Supporting Files:
• requirements.txt - Dependencies
• historical_kpi_summary.csv - KPI analysis
• run_app.bat - Application launcher"""

# Slide 14: Change Management Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Change Management Strategy"
content.text = """👥 Stakeholder Engagement:
• Site manager training programs
• Regular communication and updates
• Success story sharing across sites

📚 Knowledge Transfer:
• Documentation and user guides
• Hands-on workshops and training sessions
• Dedicated support team during transition

🎯 Adoption Incentives:
• Performance-based recognition programs
• Clear demonstration of benefits
• Gradual transition with safety nets"""

# Slide 15: Future Enhancements
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Future Enhancements"
content.text = """🔄 Real-time data integration
🧠 Advanced time series models (LSTM, Prophet)
📅 Multi-step ahead forecasting
🤖 Automated model retraining
📊 Interactive dashboard development
🔗 ERP system integration
📱 Mobile application for field teams
☁️ Cloud deployment for scalability"""

# Slide 16: Resource Requirements
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Resource Requirements"
content.text = """👨‍💻 Human Resources:
• 1 Data Scientist (Lead)
• 1 ML Engineer (Implementation)
• 1 Business Analyst (Requirements)
• Site coordinators (Part-time)

💻 Technical Infrastructure:
• Cloud computing resources
• Database storage and backup systems
• Monitoring and alerting tools

💵 Budget Allocation:
• Development: 40%
• Infrastructure: 30%
• Training & Support: 20%
• Contingency: 10%"""

# Slide 17: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You"
subtitle.text = "Questions & Discussion\n\nCement Demand Forecasting Project\nData Science Team"

# Save presentation
prs.save('C:\\Users\\hassl\\OneDrive\\Desktop\\Amdari\\Cement_Demand_Forecasting_Across_Multiple_Sites\\Cement_Demand_Forecasting_Strategic_Presentation.pptx')
print("PowerPoint presentation created successfully!")